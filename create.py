# Importando bibliotecas necessárias para o script.
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

# Função para substituir texto em um parágrafo. Usa um placeholder e substitui pelo novo texto.
def replace_text_in_paragraph_adjusted(paragraph, placeholder, new_text):
    # Verifique se new_text é um número (int ou float) e, se for, formate como porcentagem
    if isinstance(new_text, (int, float)):
        new_text = f"{new_text:.0%}"  # Isso irá converter 0.85 para "85.00%" por exemplo.
    
    # Caso contrário, apenas certifique-se de que é uma string
    else:
        new_text = str(new_text)
    
    full_text = ''.join(run.text for run in paragraph.runs)
    if placeholder in full_text:
        full_text = full_text.replace(placeholder, new_text)
        # Limpa os runs existentes antes de adicionar o novo texto
        while len(paragraph.runs) > 0:
            paragraph._p.remove(paragraph.runs[0]._r)
        paragraph.add_run().text = full_text
        return True
    return False


# Função para aplicar formatação especial (tamanho da fonte e negrito) a um parágrafo.
def apply_special_formatting(paragraph, font_size, bold=True):
    for run in paragraph.runs:
        run.font.size = font_size
        run.font.bold = bold

# Função para aplicar a cor branca ao texto de um parágrafo.
def apply_white_color_to_text(paragraph):
    for run in paragraph.runs:
        run.font.color.rgb = RGBColor(255, 255, 255)

# Função para adicionar ações numeradas com formatação específica a um quadro de texto.
def add_numbered_actions_with_formatting(text_frame, actions, font_size=Pt(14)):
    text_frame.clear()
    for i, action in enumerate(actions, start=1):
        p = text_frame.add_paragraph()
        p.text = f"{i}. {action}"
        p.font.size = font_size
        p.font.color.rgb = RGBColor(255, 255, 255)

# Função principal para atualizar a apresentação PowerPoint com base nos dados de uma linha específica.
def update_presentation_with_white_text(ppt_model_path, row):
    ppt = Presentation(ppt_model_path)
    nome = row['Nome']
    nivel = row['Nível']
    acoes = [row[f'Ação {i}'] for i in range(1, 7) if not pd.isna(row[f'Ação {i}'])]

    for slide in ppt.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    if replace_text_in_paragraph_adjusted(paragraph, "[nome]", nome):
                        apply_special_formatting(paragraph, Pt(28), bold=True)
                        apply_white_color_to_text(paragraph)
                    if replace_text_in_paragraph_adjusted(paragraph, "[nivel]", nivel):
                        apply_special_formatting(paragraph, Pt(18), bold=True)
                        apply_white_color_to_text(paragraph)
                    if "[acoes]" in ''.join(run.text for run in paragraph.runs):
                        add_numbered_actions_with_formatting(shape.text_frame, acoes)
                        for p in shape.text_frame.paragraphs:
                            apply_white_color_to_text(p)

    return ppt

# Definindo os caminhos para o modelo do PowerPoint e para a planilha de dados.
ppt_model_path = '/content/modelo.pptx'  # Caminho para o arquivo modelo do PowerPoint.
excel_data_path = '/content/template.xlsx'  # Caminho para a planilha do Excel.

# Carregando os dados da planilha de Excel.
df_mass_upload = pd.read_excel(excel_data_path)

# Processando cada linha da planilha para atualizar e salvar o PowerPoint correspondente.
for index, row in df_mass_upload.iterrows():
    # Gera a apresentação atualizada usando os dados da linha atual.
    updated_ppt = update_presentation_with_white_text(ppt_model_path, row)
    
    # Define o caminho onde o PowerPoint atualizado será salvo.
    updated_ppt_path = f'{row["Nome do Arquivo"]}'
    
    # Salva o arquivo PowerPoint atualizado no caminho especificado.
    updated_ppt.save(updated_ppt_path)
